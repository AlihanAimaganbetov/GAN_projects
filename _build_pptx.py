"""
Строит prezentaciya_1_2.pptx по требованиям PDF.
Структура: 10 слайдов: Title -> Актуальность -> Архитектура GAN ->
Датасет -> Обучение -> Результаты -> Сравнение -> Mode collapse ->
Выводы -> Спасибо.
"""
import os, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = r"C:\Users\userus\Downloads\esp"
OUT = os.path.join(ROOT, "prezentaciya_1_2.pptx")

# 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# палитра
NAVY = RGBColor(0x0F, 0x1B, 0x3D)
ACCENT = RGBColor(0xFF, 0x6B, 0x35)
LIGHT_BG = RGBColor(0xF6, 0xF7, 0xFB)
GREY = RGBColor(0x55, 0x5B, 0x6E)
DARK = RGBColor(0x1A, 0x1F, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BLANK_LAYOUT = prs.slide_layouts[6]   # Blank


def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def add_rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def slide_header(slide, num, total, section, subtitle=None):
    # верхняя полоса
    add_rect(slide, 0, 0, W, Inches(0.65), NAVY)
    add_text(slide, Inches(0.5), Inches(0.13), Inches(8), Inches(0.4),
             section, size=14, bold=True, color=WHITE)
    add_text(slide, W - Inches(2.5), Inches(0.13), Inches(2), Inches(0.4),
             f"{num} / {total}", size=14, color=WHITE, align=PP_ALIGN.RIGHT)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.5),
                 subtitle, size=28, bold=True, color=NAVY)


def bullet_block(slide, x, y, w, h, items, *, size=16, lh=0.45):
    """items = [(label, text), ...] либо [str, ...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        if isinstance(it, tuple):
            label, body = it
            r1 = p.add_run(); r1.text = f"●  {label}: "
            r1.font.bold = True; r1.font.size = Pt(size)
            r1.font.color.rgb = ACCENT; r1.font.name = "Calibri"
            r2 = p.add_run(); r2.text = body
            r2.font.size = Pt(size); r2.font.color.rgb = DARK
            r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = "●  " + it
            r.font.size = Pt(size); r.font.color.rgb = DARK
            r.font.name = "Calibri"
    return tb


def kpi_card(slide, x, y, w, h, big, label):
    add_rect(slide, x, y, w, h, LIGHT_BG)
    add_text(slide, x, y + Inches(0.25), w, Inches(0.8),
             big, size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.0), w, Inches(0.5),
             label, size=12, color=GREY, align=PP_ALIGN.CENTER)


TOTAL = 10
SECTION = "Финальный проект · GAN · 2026"

# ─────────── СЛАЙД 1 — Title ───────────
s = add_slide()
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, Inches(2.7), W, Inches(0.05), ACCENT)
add_text(s, Inches(0.8), Inches(1.4), Inches(12), Inches(0.5),
         "ФИНАЛЬНЫЙ ПРОЕКТ · GAN · 2026",
         size=18, bold=True, color=ACCENT)
add_text(s, Inches(0.8), Inches(2.9), Inches(12), Inches(2.2),
         "WGAN-GP для генерации синтетических\n"
         "сейсмо-вибрационных сигналов",
         size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(5.0), Inches(12), Inches(0.5),
         "1D-свёрточный Generative Adversarial Network на собственном датасете IMU",
         size=20, color=RGBColor(0xCB, 0xD2, 0xE0))
add_text(s, Inches(0.8), Inches(6.4), Inches(12), Inches(0.5),
         "Выполнил: студент Аймаганбетов · Дисциплина: Generative Adversarial Networks",
         size=14, color=RGBColor(0x9B, 0xA3, 0xB5))

# ─────────── СЛАЙД 2 — Актуальность ───────────
s = add_slide()
slide_header(s, 2, TOTAL, SECTION, "Актуальность темы")
add_text(s, Inches(0.5), Inches(1.7), Inches(12), Inches(0.6),
         "Зачем синтетические вибро-сигналы для сейсмо-мониторинга",
         size=20, bold=True, color=NAVY)

bullet_block(s, Inches(0.5), Inches(2.5), Inches(7.5), Inches(4),
             [("Дефицит данных",
               "реальные вибрационные события редки — собрано 2.66M отсчётов, "
               "но только 0.72% содержат полезный сигнал"),
              ("Дисбаланс классов",
               "детектор движений учится преимущественно на «тишине»"),
              ("Расширение датасета",
               "синтетические окна — путь к аугментации обучающей выборки "
               "для CNN-классификатора события"),
              ("Тестирование детектора",
               "контролируемая генерация всплесков нужной амплитуды "
               "для проверки порогов и устойчивости")],
             size=15)

kpi_card(s, Inches(8.7), Inches(2.5), Inches(2.0), Inches(1.5),
         "0.72%", "доля активных отсчётов")
kpi_card(s, Inches(10.9), Inches(2.5), Inches(2.0), Inches(1.5),
         "506", "вибро-окон в датасете")
kpi_card(s, Inches(8.7), Inches(4.2), Inches(2.0), Inches(1.5),
         "50 Hz", "частота дискретизации")
kpi_card(s, Inches(10.9), Inches(2.5) + Inches(1.7), Inches(2.0), Inches(1.5),
         "32", "отсчёта в окне (0.64 с)")

# ─────────── СЛАЙД 3 — Архитектура GAN ───────────
s = add_slide()
slide_header(s, 3, TOTAL, SECTION, "Архитектура GAN")
add_text(s, Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
         "WGAN-GP с 1D-свёрточным генератором и критиком",
         size=18, bold=True, color=NAVY)

# Generator
add_rect(s, Inches(0.5), Inches(2.3), Inches(6), Inches(4.7), LIGHT_BG)
add_text(s, Inches(0.7), Inches(2.45), Inches(5.6), Inches(0.5),
         "GENERATOR  G   ·  27 764 параметров",
         size=16, bold=True, color=ACCENT)
bullet_block(s, Inches(0.7), Inches(3.0), Inches(5.6), Inches(4),
             ["Вход: z ∈ ℝ³² (стандартный нормальный шум)",
              "Linear → reshape (64, 4)",
              "ConvTranspose1D 64→48, stride 2  → длина 8",
              "ConvTranspose1D 48→32, stride 2  → длина 16",
              "ConvTranspose1D 32→4,  stride 2  → длина 32",
              "BatchNorm + ReLU; линейный выход",
              "Результат: (4, 32) — x, y, z, extra @ 50 Гц"],
             size=13)

# Critic
add_rect(s, Inches(6.8), Inches(2.3), Inches(6), Inches(4.7), LIGHT_BG)
add_text(s, Inches(7), Inches(2.45), Inches(5.6), Inches(0.5),
         "CRITIC  D   ·  45 281 параметров",
         size=16, bold=True, color=ACCENT)
bullet_block(s, Inches(7), Inches(3.0), Inches(5.6), Inches(4),
             ["Вход: (4, 32) — реальное или сгенерированное окно",
              "Conv1D 4→32, stride 2 + LayerNorm + LeakyReLU",
              "Conv1D 32→64, stride 2 + LayerNorm + LeakyReLU",
              "Conv1D 64→128, stride 2 + LayerNorm + LeakyReLU",
              "Linear → 1 (скалярная оценка, без sigmoid)",
              "L = − E[D(real)] + E[D(fake)] + λ·GP",
              "GP = E[(‖∇D(x̂)‖₂ − 1)²], λ = 10"],
             size=13)

# ─────────── СЛАЙД 4 — Используемый датасет ───────────
s = add_slide()
slide_header(s, 4, TOTAL, SECTION, "Используемый датасет")
add_text(s, Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
         "Собственный датасет — IMU поток с ESP32 по Bluetooth",
         size=18, bold=True, color=NAVY)

bullet_block(s, Inches(0.5), Inches(2.3), Inches(8), Inches(5),
             [("Источник",
               "47 текстовых лог-файлов IMU с акселерометра, 2.66 млн отсчётов"),
              ("Каналы",
               "accel-X, accel-Y, accel-Z (g) + shake_force (RMS-метрика)"),
              ("Семплирование",
               "50 Гц, шаг 20 мс между отсчётами"),
              ("Сегментация",
               "детектор активности: extra > 0.05 g, склейка с зазором 0.5 c"),
              ("Окна",
               "длина 32 отсчёта (0.64 с), шаг 4 — итого 506 окон"),
              ("Препроцессинг",
               "вычитание DC по каналам (убрана гравитация), z-score нормализация"),
              ("Train / Val split",
               "378 / 128 (≈ 75% / 25%)")],
             size=14)

# справа коробочка с пик-extra distribution
add_rect(s, Inches(9), Inches(2.3), Inches(3.8), Inches(4.6), LIGHT_BG)
add_text(s, Inches(9.2), Inches(2.45), Inches(3.5), Inches(0.4),
         "ПИК |extra| ПО ОКНАМ",
         size=13, bold=True, color=NAVY)
bullet_block(s, Inches(9.2), Inches(2.95), Inches(3.5), Inches(4),
             [("min", "0.057 g"), ("median", "0.484 g"),
              ("max", "5.436 g"), ("std в окне", "0.23 (real) "),
              ("Источник", "разные сессии 03–04 / 2026")],
             size=12)

# ─────────── СЛАЙД 5 — Процесс обучения ───────────
s = add_slide()
slide_header(s, 5, TOTAL, SECTION, "Процесс обучения")
add_text(s, Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
         "Adversarial training: критик в N=3 шага против генератора",
         size=18, bold=True, color=NAVY)

# гиперпараметры
add_text(s, Inches(0.5), Inches(2.3), Inches(6), Inches(0.5),
         "ГИПЕРПАРАМЕТРЫ", size=14, bold=True, color=ACCENT)
bullet_block(s, Inches(0.5), Inches(2.8), Inches(6), Inches(4),
             [("Optimizer", "Adam(β1=0.5, β2=0.9), lr=1e-4"),
              ("Batch", "64"),
              ("Epochs", "80"),
              ("N_critic", "3 (критик / 1 обновление G)"),
              ("Latent dim z", "32"),
              ("λ Gradient Penalty", "10.0"),
              ("Train time на CPU", "~19 секунд"),
              ("Hardware", "Python 3.11 + PyTorch 2.4 (CPU only)")],
             size=14)

# журнал обучения
add_text(s, Inches(7), Inches(2.3), Inches(6), Inches(0.5),
         "ЛОГ W-DISTANCE", size=14, bold=True, color=ACCENT)
add_rect(s, Inches(7), Inches(2.8), Inches(5.8), Inches(4), LIGHT_BG)
log_text = ("ep   1 | D= 1.20  G= 0.34  | W-dist = 0.44\n"
            "ep  10 | D=−6.50  G= 3.48  | W-dist = 7.37\n"
            "ep  25 | D=−5.61  G= 3.84  | W-dist = 6.30\n"
            "ep  40 | D=−4.30  G= 6.53  | W-dist = 4.93\n"
            "ep  55 | D=−3.61  G= 9.99  | W-dist = 4.19\n"
            "ep  70 | D=−3.87  G=14.05  | W-dist = 4.55\n"
            "ep  80 | D=−3.86  G=15.54  | W-dist = 4.49")
add_text(s, Inches(7.2), Inches(3.0), Inches(5.5), Inches(3.8),
         log_text, size=12, color=DARK, font="Consolas")

# ─────────── СЛАЙД 6 — Результаты генерации ───────────
s = add_slide()
slide_header(s, 6, TOTAL, SECTION, "Результаты генерации")
add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.5),
         "Сравнение реальных и синтетических окон: формы и спектры",
         size=18, bold=True, color=NAVY)

img_path = os.path.join(ROOT, "_compare_v2.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.6), Inches(2.1),
                         width=Inches(9.2))
# легенда / комментарий справа
add_rect(s, Inches(10.1), Inches(2.1), Inches(2.8), Inches(4.7), LIGHT_BG)
add_text(s, Inches(10.25), Inches(2.25), Inches(2.6), Inches(0.4),
         "ЧИТАЕМ ГРАФИК", size=12, bold=True, color=NAVY)
bullet_block(s, Inches(10.25), Inches(2.7), Inches(2.6), Inches(4),
             ["Строка 1 — real окна",
              "Строка 2 — fake окна",
              "Строка 3 — средний |FFT|",
              "x/y/z/extra — каналы",
              "Y одинаков real/fake → амплитуды сравнимы"],
             size=11)

# ─────────── СЛАЙД 7 — Сравнение real/fake ───────────
s = add_slide()
slide_header(s, 7, TOTAL, SECTION, "Сравнение real vs fake")
add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.5),
         "Метрики покрытия и распределения по каналам",
         size=18, bold=True, color=NAVY)

# таблица метрик
table_x, table_y = Inches(0.5), Inches(2.2)
rows, cols = 5, 6
tbl_shape = s.shapes.add_table(rows, cols, table_x, table_y, Inches(7.5), Inches(2.2))
tbl = tbl_shape.table
headers = ["Канал", "real μ", "fake μ", "real σ", "fake σ", "W₁"]
data = [
    ["x",     "0.0000", "−0.0070", "0.2271", "0.1139", "0.0384"],
    ["y",    "−0.0000", "−0.0003", "0.1742", "0.0847", "0.0311"],
    ["z",     "0.0000",  "0.0028", "0.2526", "0.1391", "0.0399"],
    ["extra","−0.0000",  "0.0140", "0.3708", "0.1882", "0.0777"],
]
for j, h in enumerate(headers):
    c = tbl.cell(0, j)
    c.text = h
    p = c.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for r_ in p.runs:
        r_.font.bold = True; r_.font.size = Pt(13)
        r_.font.color.rgb = WHITE
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
for i, row in enumerate(data, 1):
    for j, val in enumerate(row):
        c = tbl.cell(i, j); c.text = val
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for r_ in p.runs:
            r_.font.size = Pt(13)
            r_.font.color.rgb = DARK
            if j == 0: r_.font.bold = True

# гистограммы справа
img_path = os.path.join(ROOT, "_hist_v2.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.5), Inches(4.7),
                         width=Inches(12.3))

# вывод
add_text(s, Inches(8.4), Inches(2.2), Inches(4.5), Inches(2.2),
         "ИНТЕРПРЕТАЦИЯ",
         size=13, bold=True, color=ACCENT)
bullet_block(s, Inches(8.4), Inches(2.55), Inches(4.5), Inches(2.4),
             ["Средние совпадают (≈ 0): DC-вычитание ок",
              "σ fake ≈ ½ σ real — крайние пики недо-учтены",
              "W₁ < 0.08 по всем каналам — слабая дивергенция",
              "Спектры x/y/z совпадают, extra — расхождение"],
             size=11)

# ─────────── СЛАЙД 8 — Mode collapse и стабилизация ───────────
s = add_slide()
slide_header(s, 8, TOTAL, SECTION, "Борьба с mode collapse")
add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.5),
         "Эксперимент v1 (BCE-GAN) → v2 (WGAN-GP): что изменилось",
         size=18, bold=True, color=NAVY)

# колонка v1
add_rect(s, Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.9), LIGHT_BG)
add_text(s, Inches(0.7), Inches(2.35), Inches(5.5), Inches(0.45),
         "ИТЕРАЦИЯ v1 · BCE-GAN  ✗",
         size=15, bold=True, color=RGBColor(0xB0, 0x1E, 0x1E))
bullet_block(s, Inches(0.7), Inches(2.85), Inches(5.5), Inches(4.5),
             ["Порог extra > 0.015 → 72 337 окон",
              "BCE loss, без gradient penalty",
              "Генератор схлопнулся к фоновому шуму",
              "FFT-спектр fake → ложные пики",
              "Гистограммы амплитуд: fake — узкий пик у нуля",
              "Дискриминатор: acc_real ≈ 0.4 — путается"],
             size=12)

# колонка v2
add_rect(s, Inches(6.5), Inches(2.2), Inches(6.3), Inches(4.9), LIGHT_BG)
add_text(s, Inches(6.7), Inches(2.35), Inches(6), Inches(0.45),
         "ИТЕРАЦИЯ v2 · WGAN-GP  ✓",
         size=15, bold=True, color=RGBColor(0x16, 0x82, 0x4D))
bullet_block(s, Inches(6.7), Inches(2.85), Inches(6), Inches(4.5),
             ["Жёсткий порог extra > 0.05 → 506 окон, но реально вибрационных",
              "DC-вычитание по окну — модель учит форму, не уровни",
              "Wasserstein-1 + Gradient Penalty (λ=10)",
              "LayerNorm в критике (не BatchNorm — нарушает GP)",
              "N_critic = 3 — критик обновляется чаще G",
              "BatchNorm в генераторе сохранён — стабилизатор G",
              "W-distance падает с 7.4 до 4.5 — есть прогресс",
              "Спектры x/y/z совпадают с real, нет ложных пиков"],
             size=12)

# ─────────── СЛАЙД 9 — Выводы и future work ───────────
s = add_slide()
slide_header(s, 9, TOTAL, SECTION, "Выводы и future work")

add_text(s, Inches(0.5), Inches(1.55), Inches(6), Inches(0.5),
         "ДОСТИГНУТО", size=18, bold=True, color=ACCENT)
bullet_block(s, Inches(0.5), Inches(2.1), Inches(6.2), Inches(5),
             ["Полностью реализован adversarial training pipeline на PyTorch",
              "Generator + Discriminator → WGAN-GP с λ_GP = 10",
              "Использованы методы устойчивости: BatchNorm в G, LayerNorm в D, GP, "
              "DC-вычитание, осторожный betas Adam",
              "Достигнуто совпадение спектральной структуры реальных и синтетических "
              "окон по x, y, z",
              "Артефакты обучения: _gan_v2.pt (298 КБ), _samples_v2.npy, generate.py",
              "Скрипт generate.py — CLI-генератор окон в CSV для дальнейшего "
              "использования в проекте"],
             size=13)

add_text(s, Inches(7), Inches(1.55), Inches(6), Inches(0.5),
         "FUTURE WORK", size=18, bold=True, color=ACCENT)
bullet_block(s, Inches(7), Inches(2.1), Inches(6.2), Inches(5),
             ["Conditional GAN: метка пиковой амплитуды (массив amp_tr уже "
              "в датасете) — контроль силы вибрации",
              "Увеличить датасет до ~5000 реальных вибро-окон — "
              "снять занижение σ для extra-канала",
              "Заменить single-scale CNN на multi-scale "
              "(чтобы поймать одновременно быстрые транзиенты и медленные тренды)",
              "Spectral consistency loss — прямое выравнивание FFT-спектров",
              "EMA весов генератора — стандартный приём WGAN для качества образцов",
              "Интеграция с CNN-классификатором: тестирование детектора "
              "на синтетических окнах с заданной амплитудой"],
             size=13)

# ─────────── СЛАЙД 10 — Спасибо ───────────
s = add_slide()
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, Inches(5), Inches(3.4), Inches(3.3), Inches(0.07), ACCENT)
add_text(s, 0, Inches(2.2), W, Inches(1.2),
         "СПАСИБО", size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(3.7), W, Inches(0.7),
         "за внимание", size=28, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(4.9), W, Inches(0.5),
         "Готов ответить на ваши вопросы",
         size=16, color=RGBColor(0xCB, 0xD2, 0xE0), align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(6.7), W, Inches(0.4),
         "WGAN-GP для генерации сейсмо-вибрационных сигналов  ·  Финальный проект",
         size=12, color=RGBColor(0x9B, 0xA3, 0xB5), align=PP_ALIGN.CENTER)

# сохраняем
prs.save(OUT)
print(f"Слайдов: {len(prs.slides)}")
print(f"Сохранено: {OUT}")
print(f"Размер: {os.path.getsize(OUT)/1024:.1f} КБ")
