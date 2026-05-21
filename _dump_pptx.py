"""Извлекает текст из всех слайдов pptx."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation

p = Presentation(r"C:\Users\userus\Downloads\esp\prezentaciya_1_1.pptx")
print(f"Слайдов: {len(p.slides)}\n")
for i, s in enumerate(p.slides, 1):
    print(f"=== Слайд {i} ===")
    for shape in s.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs).strip()
                if txt:
                    print(f"  {txt}")
        elif shape.shape_type == 13:  # picture
            print(f"  [картинка: {shape.name}]")
        elif shape.has_table:
            print(f"  [таблица: {shape.name}]")
    print()
